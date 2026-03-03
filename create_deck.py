#!/usr/bin/env python3
"""
Engrm Pitch Deck - a16z speedrun
Premium yellow/black design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# Yellow/Black palette (a16z style)
YELLOW = RGBColor(255, 219, 21)
YELLOW_DARK = RGBColor(244, 144, 12)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(59, 130, 246)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_yellow_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = YELLOW
    bg.line.fill.background()

def add_black_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, fill_color=DARK_GRAY, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.fill.background()
    return card

def txt(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align

def mono(slide, left, top, width, height, text, size=11, color=GRAY, align=PP_ALIGN.LEFT):
    txt(slide, left, top, width, height, text, size, color, False, align, "Consolas")

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)
add_shape(slide, 10, 1.5, 6, 6, DARK_GRAY)  # Brain orb

txt(slide, 1.2, 3, 9, 1.6, "Memory for\nAI agents", size=72, color=BLACK, bold=True)
txt(slide, 1.2, 5.5, 8, 1, "Persistent memory that compounds over time.\nLike a human brain, not a chat log.", size=22, color=DARK_GRAY)
txt(slide, 1.2, 7.2, 6, 0.5, "engrm.xyz", size=18, color=BLACK, bold=True, font="Consolas")
mono(slide, 14.5, 8.2, 1.2, 0.4, "01", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1.2, "Every AI agent\nwakes up with amnesia", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 4, 8, 1.5, "ChatGPT forgets you exist between sessions.\nCustom agents lose context every restart.\nRAG retrieves documents, not relationships.", size=22, color=DARK_GRAY)

# Stats
add_card(slide, 10, 2, 4.5, 1.8, BLACK)
txt(slide, 10.3, 2.2, 4, 0.9, "0", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 3, 4, 0.5, "MEMORIES RETAINED\nBETWEEN SESSIONS", size=10, color=GRAY)

add_card(slide, 10, 4, 4.5, 1.8, BLACK)
txt(slide, 10.3, 4.2, 4, 0.9, "100%", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 5, 4, 0.5, "CONTEXT LOST\nON RESTART", size=10, color=GRAY)

add_card(slide, 10, 6, 4.5, 1.8, BLACK)
txt(slide, 10.3, 6.2, 4, 0.9, "$47B", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 7, 4, 0.5, "AI AGENT MARKET\nBY 2030", size=10, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "02", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 3: Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1, "A brain, not a database", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 3, 8, 0.8, "Memory that works like human cognition.", size=22, color=DARK_GRAY)

# Three pillars
pillars = [
    ("🧠", "Biological decay", "Memories fade naturally.\nImportant ones strengthen."),
    ("🔗", "Associative linking", "Related memories connect.\nContext compounds."),
    ("🔒", "Privacy-first", "Client-side encryption.\nYour data stays yours."),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 1.2 + i * 5
    add_card(slide, x, 4, 4.5, 3.2, BLACK)
    txt(slide, x + 0.2, 4.3, 4, 0.6, icon, size=36, color=WHITE)
    txt(slide, x + 0.2, 5, 4, 0.5, title, size=20, color=YELLOW, bold=True)
    txt(slide, x + 0.2, 5.6, 4, 1.4, desc, size=14, color=WHITE)

mono(slide, 14.5, 8.2, 1.2, 0.4, "03", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 4: How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 12, 1, "How it works", size=48, color=WHITE, bold=True)

# Flow
steps = [
    ("Store", "Agent saves memory\nvia API"),
    ("Embed", "Vectorize +\nlink to graph"),
    ("Decay", "Strengthen or\nfade over time"),
    ("Recall", "Semantic search +\ncontext inject"),
]
for i, (title, desc) in enumerate(steps):
    x = 1.2 + i * 3.7
    add_card(slide, x, 3.5, 3.2, 2.5, DARK_GRAY, YELLOW)
    txt(slide, x + 0.2, 3.8, 2.8, 0.5, title, size=20, color=YELLOW, bold=True)
    txt(slide, x + 0.2, 4.4, 2.8, 1.5, desc, size=14, color=WHITE)
    if i < 3:
        txt(slide, x + 3.3, 4.3, 0.5, 0.5, "→", size=24, color=YELLOW)

# Code snippet
mono(slide, 1.2, 6.8, 14, 1.5, """# Store a memory
engrm.store("User prefers dark mode and minimalist design")

# Recall relevant context  
memories = engrm.recall("What design style does the user like?")""", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "04", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 5: Memory Lifecycle
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Memory lifecycle", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.6, "Biological decay curves. Not FIFO deletion.", size=20, color=DARK_GRAY)

# Timeline
add_card(slide, 1.2, 3.8, 13.5, 3.5, BLACK)

# Day markers
for i, (day, event, strength) in enumerate([
    ("Day 1", "Memory created", "100%"),
    ("Day 7", "No access, begins decay", "85%"),
    ("Day 14", "Mentioned once, reinforced", "95%"),
    ("Day 30", "Accessed in conversation", "90%"),
    ("Day 90", "Low access, archived", "40%"),
]):
    x = 1.5 + i * 2.6
    txt(slide, x, 4.1, 2.4, 0.4, day, size=12, color=YELLOW, bold=True)
    txt(slide, x, 4.5, 2.4, 0.8, event, size=11, color=WHITE)
    txt(slide, x, 5.5, 2.4, 0.4, strength, size=18, color=YELLOW, bold=True)

txt(slide, 1.5, 6.5, 12, 0.6, "Identity memories (\"my name is...\") → Never deleted, 365-day half-life", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "05", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 6: Integrations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Works with everything", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.6, "Any LLM. Any framework. 5 minutes to integrate.", size=20, color=DARK_GRAY)

# Integration cards
integrations = [
    ("REST API", "Standard HTTP endpoints.\nJSON in, JSON out."),
    ("MCP Server", "Model Context Protocol.\nWorks with Claude, etc."),
    ("Python SDK", "pip install engrm\nOne import, done."),
]
for i, (title, desc) in enumerate(integrations):
    x = 1.2 + i * 5
    add_card(slide, x, 3.8, 4.5, 2.2, BLACK)
    txt(slide, x + 0.3, 4.1, 4, 0.5, title, size=18, color=YELLOW, bold=True)
    txt(slide, x + 0.3, 4.7, 4, 1.2, desc, size=14, color=WHITE)

# Compatible with
txt(slide, 1.2, 6.5, 14, 0.5, "Compatible with: OpenAI · Anthropic · Llama · Mistral · Any LLM", size=16, color=DARK_GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "06", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 7: Privacy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Privacy-first", size=48, color=WHITE, bold=True)
txt(slide, 1.2, 2.6, 8, 0.8, "Client-side encryption option.\nWe can't read your memories even if we wanted to.", size=20, color=GRAY)

# Encryption flow
add_card(slide, 1.2, 4, 4, 2.5, DARK_GRAY, YELLOW)
txt(slide, 1.5, 4.3, 3.5, 0.5, "Your device", size=16, color=YELLOW, bold=True)
txt(slide, 1.5, 4.9, 3.5, 1.5, "AES-256-GCM\nencryption\nbefore upload", size=14, color=WHITE)

txt(slide, 5.5, 5, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 6.5, 4, 4, 2.5, DARK_GRAY)
txt(slide, 6.8, 4.3, 3.5, 0.5, "Our servers", size=16, color=GRAY, bold=True)
txt(slide, 6.8, 4.9, 3.5, 1.5, "Store ciphertext\n+ encrypted\nvectors", size=14, color=GRAY)

txt(slide, 10.8, 5, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 11.5, 4, 4, 2.5, DARK_GRAY, YELLOW)
txt(slide, 11.8, 4.3, 3.5, 0.5, "Your device", size=16, color=YELLOW, bold=True)
txt(slide, 11.8, 4.9, 3.5, 1.5, "Decrypt on\nretrieval\nlocally", size=14, color=WHITE)

txt(slide, 1.2, 7.2, 12, 0.5, "Optional: Permanent storage on Arweave (premium add-on)", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "07", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 8: Market
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Why now", size=48, color=BLACK, bold=True)

points = [
    "Agents are shipping — ChatGPT Plugins, Claude MCP, custom bots everywhere",
    "Memory is the missing layer — RAG solves documents, not relationships",
    "Privacy matters — users want AI that remembers them, not corporations",
    "Foundation models commoditizing — memory becomes the moat",
]
for i, point in enumerate(points):
    txt(slide, 1.2, 2.8 + i * 1.2, 13, 0.9, f"→  {point}", size=20, color=BLACK)

# TAM
add_card(slide, 10, 6.5, 4.5, 1.5, BLACK)
txt(slide, 10.3, 6.7, 4, 0.7, "$47B", size=36, color=YELLOW, bold=True)
mono(slide, 10.3, 7.4, 4, 0.4, "AI AGENT TAM 2030", size=10, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "08", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 9: Pricing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Simple pricing", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.5, "Start free. Scale with usage.", size=20, color=DARK_GRAY)

# Tiers
tiers = [
    ("FREE", "$0", "5K memories/mo\n100MB storage\nCommunity support"),
    ("PRO", "$29", "50K memories/mo\n1GB storage\nPriority support"),
    ("ENTERPRISE", "Custom", "Unlimited\nOn-prem option\nSLA + support"),
]
for i, (tier, price, features) in enumerate(tiers):
    x = 1.2 + i * 5
    is_pro = tier == "PRO"
    add_card(slide, x, 3.5, 4.5, 4, BLACK if is_pro else DARK_GRAY)
    txt(slide, x + 0.3, 3.8, 4, 0.4, tier, size=12, color=YELLOW if is_pro else GRAY, bold=True, font="Consolas")
    txt(slide, x + 0.3, 4.3, 4, 0.7, price, size=44, color=YELLOW if is_pro else WHITE, bold=True)
    txt(slide, x + 0.3, 5.2, 4, 2, features, size=14, color=WHITE if is_pro else GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "09", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 10: CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)
add_shape(slide, 9, 2, 7, 7, DARK_GRAY)  # Brain orb

txt(slide, 1.2, 3, 8, 1.6, "Give your agents\na memory", size=64, color=WHITE, bold=True)
txt(slide, 1.2, 5.5, 7, 1, "Persistent. Private. Biological.\nThe memory layer for AI.", size=22, color=GRAY)

add_card(slide, 1.2, 7, 3.5, 0.7, YELLOW)
txt(slide, 1.4, 7.1, 3.2, 0.5, "Start building →", size=16, color=BLACK, bold=True)

txt(slide, 1.2, 8, 4, 0.4, "engrm.xyz", size=18, color=YELLOW, bold=True, font="Consolas")

mono(slide, 14.5, 8.2, 1.2, 0.4, "10", color=YELLOW, align=PP_ALIGN.RIGHT)

# Save
prs.save('/Users/clawdaddy/clawd/projects/engrm/Engrm-Pitch-Deck.pptx')
print("Created: Engrm-Pitch-Deck.pptx")
