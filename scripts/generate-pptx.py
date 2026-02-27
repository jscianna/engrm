#!/usr/bin/env python3
"""Generate MEMRY presentation as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
CYAN = RGBColor(34, 211, 238)
DARK_BG = RGBColor(9, 9, 11)
ZINC_800 = RGBColor(39, 39, 42)
ZINC_400 = RGBColor(161, 161, 170)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(239, 68, 68)
GREEN = RGBColor(34, 197, 94)
PURPLE = RGBColor(168, 85, 247)
AMBER = RGBColor(245, 158, 11)
PINK = RGBColor(236, 72, 153)
BLUE = RGBColor(59, 130, 246)

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text, top=0.5, size=44, color=WHITE):
    """Add a title to the slide."""
    left = Inches(0.75)
    width = Inches(12.5)
    height = Inches(1)
    
    shape = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color

def add_subtitle(slide, text, top=1.2, size=24, color=ZINC_400):
    """Add a subtitle."""
    left = Inches(0.75)
    width = Inches(12.5)
    height = Inches(0.6)
    
    shape = slide.shapes.add_textbox(left, Inches(top), width, height)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color

def add_text_box(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return shape

def add_bullet_list(slide, items, left, top, width, size=20, color=WHITE, bullet_color=CYAN):
    """Add a bullet list."""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4))
    tf = shape.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→ {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_before = Pt(8)

def add_logo(slide):
    """Add MEMRY logo."""
    add_text_box(slide, "MEMRY", 0.5, 0.3, 2, 0.5, size=16, color=CYAN, bold=True)

def add_slide_number(slide, num):
    """Add slide number."""
    add_text_box(slide, str(num), 12.5, 7, 0.5, 0.3, size=12, color=ZINC_400)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_background(slide, DARK_BG)
add_text_box(slide, "MEMRY", 0, 2.5, 13.333, 1.5, size=72, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, "Zero-Knowledge Memory for AI Agents", 0, 3.8, 13.333, 0.8, size=32, color=ZINC_400, align=PP_ALIGN.CENTER)
add_text_box(slide, "Permanent Storage  •  True Encryption  •  Local Embeddings  •  Biological Decay", 0, 5, 13.333, 0.6, size=16, color=CYAN, align=PP_ALIGN.CENTER)
add_slide_number(slide, 1)

# Slide 2: The Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "The Problem")

add_text_box(slide, "Current AI Memory", 1, 1.5, 5, 0.5, size=24, color=RED, bold=True)
add_bullet_list(slide, [
    "Context window limitations",
    "No persistence across sessions",
    "Centralized, readable by providers",
    "No biological memory dynamics",
    "Flat storage, no associations"
], 1, 2.1, 5, size=18, color=ZINC_400)

add_text_box(slide, "What Agents Need", 7, 1.5, 5, 0.5, size=24, color=GREEN, bold=True)
add_bullet_list(slide, [
    "Unlimited long-term memory",
    "True privacy (zero-knowledge)",
    "Intelligent recall & linking",
    "Natural forgetting (decay)",
    "Associative memory graph"
], 7, 2.1, 5, size=18, color=ZINC_400)
add_slide_number(slide, 2)

# Slide 3: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Zero-Knowledge Architecture")

# Flow boxes
flow_items = [
    ("💬", "Conversation"),
    ("🧠", "Heuristics\nscore ≥ 6.0?"),
    ("📊", "Local Embed\nFastEmbed ONNX"),
    ("🔒", "Encrypt\nAES-256-GCM"),
    ("☁️", "MEMRY\nStore + Link"),
]

x_start = 0.5
for i, (emoji, label) in enumerate(flow_items):
    x = x_start + i * 2.5
    # Box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2), Inches(2), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ZINC_800
    shape.line.color.rgb = CYAN if i == 3 else ZINC_800
    
    # Emoji
    add_text_box(slide, emoji, x, 2.1, 2, 0.5, size=28, align=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, label, x, 2.7, 2, 0.8, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    
    # Arrow
    if i < len(flow_items) - 1:
        add_text_box(slide, "→", x + 2, 2.5, 0.5, 0.5, size=24, color=CYAN)

# Bottom box
add_text_box(slide, "Server receives: Vector + Encrypted blob", 0, 4.5, 13.333, 0.4, size=18, color=CYAN, align=PP_ALIGN.CENTER)
add_text_box(slide, "Server CANNOT: Read content • See queries • Decrypt data", 0, 5, 13.333, 0.4, size=18, color=RED, align=PP_ALIGN.CENTER)
add_slide_number(slide, 3)

# Slide 4: Context Injection (How it works with LLMs)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "How Context Injection Works")
add_subtitle(slide, "Every turn, agent retrieves relevant memories and injects into prompt")

# Flow diagram
flow_steps = [
    ("1", "User Message", "What time works\nfor a call?"),
    ("2", "Embed & Query", "MEMRY returns\ntop-K memories"),
    ("3", "Decrypt Local", "Agent decrypts\nwith vault key"),
    ("4", "Inject Context", "Add to system\nprompt"),
    ("5", "LLM Call", "Response uses\nmemory context"),
]

x = 0.3
for num, title, desc in flow_steps:
    # Circle with number
    add_text_box(slide, num, x + 0.9, 2.1, 0.5, 0.4, size=20, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x, 2.6, 2.3, 0.4, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x, 3.1, 2.3, 0.8, size=11, color=ZINC_400, align=PP_ALIGN.CENTER)
    if num != "5":
        add_text_box(slide, "→", x + 2.3, 2.6, 0.3, 0.4, size=20, color=CYAN)
    x += 2.5

# Example prompt box
add_text_box(slide, "Example Augmented Prompt:", 0.75, 4.2, 4, 0.4, size=16, color=CYAN, bold=True)
prompt_example = """System: You are a helpful assistant.

## Relevant memories:
- User prefers morning meetings (preference)
- User timezone is SGT, GMT+8 (fact)
- User is busy Fridays (constraint)

User: What time works for a call?"""
add_text_box(slide, prompt_example, 0.75, 4.7, 5.5, 2.5, size=12, color=ZINC_400)

# Key points
add_text_box(slide, "Key Benefits:", 7, 4.2, 5, 0.4, size=16, color=CYAN, bold=True)
add_bullet_list(slide, [
    "Retrieval per turn — always fresh context",
    "Semantic search — meaning, not keywords",
    "Strength-weighted — important memories first",
    "Co-retrieval bonds — related memories link",
    "Decay filters — old noise fades away"
], 7, 4.7, 5, size=13, color=ZINC_400)

add_slide_number(slide, 4)

# Slide 5: Heuristic Scoring
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Intelligent Extraction (No LLM Cost)")
add_subtitle(slide, "Deterministic heuristics decide what's worth remembering")

signals = [
    ("Explicit Markers", "+3.0", '"remember this", "my name is"', CYAN),
    ("Decision", "+2.0", '"decided", "going with"', PURPLE),
    ("Correction", "+1.5", '"actually", "no wait"', AMBER),
    ("Emotional", "+1.5", 'CAPS, "amazing!!!"', PINK),
    ("Temporal", "+1.0", '"January 5th", "next Tuesday"', GREEN),
    ("Entities", "+2.0", 'Names, URLs, dates', BLUE),
]

y = 2
for signal, points, example, color in signals:
    add_text_box(slide, signal, 0.75, y, 2.5, 0.4, size=16, color=color, bold=True)
    add_text_box(slide, points, 3.5, y, 1, 0.4, size=16, color=WHITE, bold=True)
    add_text_box(slide, example, 5, y, 7, 0.4, size=14, color=ZINC_400)
    y += 0.55

add_text_box(slide, "Final Score = Base (0-7) × Type Multiplier (0.8-1.3) → Store if ≥ 6.0", 0, 6, 13.333, 0.4, size=16, color=ZINC_400, align=PP_ALIGN.CENTER)
add_slide_number(slide, 5)

# Slide 6: Memory Types & Decay
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Memory Types & Decay")
add_subtitle(slide, "Biological memory dynamics — important things persist, trivia fades")

types = [
    ("Constraint", "180 days", 1.0, RED),
    ("Identity", "120 days", 0.67, PURPLE),
    ("How-to", "120 days", 0.67, GREEN),
    ("Fact", "90 days", 0.5, BLUE),
    ("Preference", "60 days", 0.33, CYAN),
    ("Relationship", "30 days", 0.17, PINK),
    ("Event", "14 days", 0.08, AMBER),
]

y = 2
for name, days, width_pct, color in types:
    add_text_box(slide, f"● {name}", 0.75, y, 2, 0.35, size=14, color=color)
    # Bar
    bar_width = 6 * width_pct
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(y), Inches(bar_width), Inches(0.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.fill.background()
    add_text_box(slide, days, 9.5, y, 1.5, 0.35, size=12, color=ZINC_400)
    y += 0.5

add_text_box(slide, "strength = base × (0.9 ^ (days_since_access / halflife))", 0, 6, 13.333, 0.3, size=14, color=ZINC_400, align=PP_ALIGN.CENTER)
add_text_box(slide, "Access resets the clock • Repeated mentions strengthen • Below 0.3 → archive → delete", 0, 6.4, 13.333, 0.3, size=12, color=ZINC_400, align=PP_ALIGN.CENTER)
add_slide_number(slide, 6)

# Slide 7: Reinforcement
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Fire Together, Wire Together")
add_subtitle(slide, "Repeated mentions strengthen memories instead of duplicating")

# Left side - logic
add_text_box(slide, "Reinforcement Logic", 0.75, 2, 5, 0.5, size=20, color=CYAN, bold=True)
logic_text = """New memory arrives
       ↓
Find similar (cosine > 0.85)?
       ↓
  YES → REINFORCE
  • strength ↑
  • mentions ↑
  • merge entities

  NO → CREATE NEW
  • strength = 1.0
  • set halflife
  • auto-link similar"""
add_text_box(slide, logic_text, 0.75, 2.6, 5, 4, size=14, color=ZINC_400)

# Right side - frequency boost
add_text_box(slide, "Frequency Boost", 7, 2, 5, 0.5, size=20, color=CYAN, bold=True)
boosts = [
    ("1st mention", "1.0×"),
    ("2nd mention", "1.4× (+40%)"),
    ("3rd mention", "1.6× (+20%)"),
    ("4th mention", "1.8× (+20%)"),
    ("5th+ mention", "+5% each"),
    ("Maximum", "2.5× cap"),
]
y = 2.6
for label, value in boosts:
    add_text_box(slide, label, 7, y, 3, 0.35, size=14, color=WHITE)
    add_text_box(slide, value, 10.5, y, 2, 0.35, size=14, color=CYAN)
    y += 0.45

add_text_box(slide, 'Example: "Rex" mentioned 5× → 0.6 → 0.84 → 0.97 → 1.08 → 1.13', 7, 5.5, 5, 0.4, size=12, color=ZINC_400)
add_slide_number(slide, 7)

# Slide 8: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Technology Stack")

# Three columns
columns = [
    ("Storage", ["Turso — Edge SQLite", "Arweave — Permanent", "Vector Index — Similarity"]),
    ("Security", ["AES-256-GCM — Encryption", "PBKDF2 — Key derivation", "Client-side — All crypto"]),
    ("Intelligence", ["FastEmbed — Local ONNX", "Heuristics — No LLM cost", "Graph — Associations"]),
]

x = 0.75
for title, items in columns:
    add_text_box(slide, title, x, 1.5, 3.5, 0.5, size=20, color=CYAN, bold=True)
    add_bullet_list(slide, items, x, 2.1, 3.5, size=14, color=ZINC_400)
    x += 4.2

# Stats
stats = [
    ("$0", "LLM Cost"),
    ("384", "Dimensions"),
    ("~200ms", "Embed Time"),
    ("∞", "Storage"),
]
x = 1.5
for value, label in stats:
    add_text_box(slide, value, x, 5, 2.5, 0.6, size=36, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, x, 5.6, 2.5, 0.4, size=12, color=ZINC_400, align=PP_ALIGN.CENTER)
    x += 2.8
add_slide_number(slide, 8)

# Slide 9: Integration
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_logo(slide)
add_title(slide, "Easy Integration")

integrations = [
    ("MCP Server", "Claude Desktop & MCP clients", 'npm install -g @memry/mcp\n\n"memry": { "command": "memry-mcp" }'),
    ("Python SDK", "Full ZK with FastEmbed", 'pip install memry-sdk\n\nclient.store("User likes dark theme")'),
    ("REST API", "Any platform", 'POST /api/v1/memories\n{ "content": "...", "embedding": [...] }'),
]

x = 0.5
for title, desc, code in integrations:
    add_text_box(slide, title, x, 1.5, 4, 0.5, size=20, color=CYAN, bold=True)
    add_text_box(slide, desc, x, 2, 4, 0.4, size=12, color=ZINC_400)
    add_text_box(slide, code, x, 2.5, 4, 3, size=11, color=WHITE)
    x += 4.2
add_slide_number(slide, 9)

# Slide 10: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_background(slide, DARK_BG)
add_text_box(slide, "Give Your Agent", 0, 2, 13.333, 0.8, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, "Real Memory", 0, 2.8, 13.333, 0.8, size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

pillars = [
    ("🔒", "Zero-Knowledge", "We can't read your data"),
    ("🧠", "Biological Dynamics", "Decay, reinforcement, linking"),
    ("💰", "Zero LLM Cost", "Heuristics for extraction"),
    ("♾️", "Permanent", "Arweave storage forever"),
]

x = 1.5
for emoji, title, desc in pillars:
    add_text_box(slide, emoji, x, 4.2, 2.5, 0.6, size=36, align=PP_ALIGN.CENTER)
    add_text_box(slide, title, x, 4.9, 2.5, 0.4, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, x, 5.4, 2.5, 0.4, size=11, color=ZINC_400, align=PP_ALIGN.CENTER)
    x += 2.7

add_text_box(slide, "memry.ai", 0, 6.5, 13.333, 0.5, size=24, color=CYAN, align=PP_ALIGN.CENTER)
add_slide_number(slide, 10)

# Save
output_path = os.path.expanduser("~/clawd/projects/memry/MEMRY-Presentation.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
