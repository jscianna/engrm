#!/usr/bin/env python3
"""
Engrm Pitch Deck v3 - Updated with all March 2026 features
Key changes:
- "The cognitive layer" positioning
- Server-side encryption (not client-side, but still "user & agent only")
- Tiered importance (critical/working/high/normal)
- Memory types & classification
- Reflect, consolidate, decay
- Python SDK, Memory Browser
- LongMemEval benchmark integration
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Color palette
YELLOW = RGBColor(255, 219, 21)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(34, 211, 238)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_yellow_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = YELLOW
    bg.line.fill.background()

def add_black_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, fill_color=DARK_GRAY, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.fill.background()
    return card

def txt(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align

def mono(slide, left, top, width, height, text, size=11, color=GRAY, align=PP_ALIGN.LEFT):
    txt(slide, left, top, width, height, text, size, color, False, align, "Consolas")

# ============================================================
# SLIDE 1: Title - "The Cognitive Layer"
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)
add_shape(slide, 10, 1.5, 6, 6, DARK_GRAY)

txt(slide, 1.2, 2.5, 9, 1.6, "The cognitive\nlayer", size=72, color=BLACK, bold=True)
txt(slide, 1.2, 5.3, 8, 1.2, "Memory infrastructure for AI agents.\nEvery conversation builds on the last.", size=24, color=DARK_GRAY)
txt(slide, 1.2, 7, 6, 0.5, "Persistent Intelligence", size=14, color=GRAY, font="Consolas")
txt(slide, 1.2, 7.5, 6, 0.5, "engrm.xyz", size=18, color=BLACK, bold=True, font="Consolas")
mono(slide, 14.5, 8.2, 1.2, 0.4, "01", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1.2, "Every AI agent\nwakes up with amnesia", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 4, 8, 1.5, "ChatGPT forgets you between sessions.\nCustom agents lose context on restart.\nRAG retrieves documents, not relationships.", size=22, color=DARK_GRAY)

add_card(slide, 10, 2, 4.5, 1.8, BLACK)
txt(slide, 10.3, 2.2, 4, 0.9, "0", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 3, 4, 0.5, "MEMORIES RETAINED\nBETWEEN SESSIONS", size=10, color=GRAY)

add_card(slide, 10, 4, 4.5, 1.8, BLACK)
txt(slide, 10.3, 4.2, 4, 0.9, "$7.5M", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 5, 4, 0.5, "COGNEE SEED ROUND\nVALIDATES CATEGORY", size=10, color=GRAY)

add_card(slide, 10, 6, 4.5, 1.8, BLACK)
txt(slide, 10.3, 6.2, 4, 0.9, "$47B", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 7, 4, 0.5, "AI AGENT MARKET\nBY 2030", size=10, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "02", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 3: Solution - A Brain Not a Database
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1, "A brain, not a database", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 3, 8, 0.8, "Memory that works like human cognition — and improves itself.", size=22, color=DARK_GRAY)

pillars = [
    ("🧠", "Intelligent recall", "Memory types: identity, preference,\nbelief, decision, fact. Context-aware."),
    ("🔗", "Self-improving", "Reflect, consolidate, decay.\nMemories that help get reinforced."),
    ("🔒", "User & Agent only", "Server-side AES-256-GCM.\nNot even we can read your data."),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 1.2 + i * 5
    add_card(slide, x, 4, 4.5, 3.2, BLACK)
    txt(slide, x + 0.2, 4.3, 4, 0.6, icon, size=36, color=WHITE)
    txt(slide, x + 0.2, 5, 4, 0.5, title, size=20, color=YELLOW, bold=True)
    txt(slide, x + 0.2, 5.6, 4, 1.4, desc, size=14, color=WHITE)

mono(slide, 14.5, 8.2, 1.2, 0.4, "03", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 4: How It Works - Simple API
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 12, 1, "Simple API. Powerful system.", size=48, color=WHITE, bold=True)

add_card(slide, 1.2, 3, 4.5, 4.5, DARK_GRAY, YELLOW)
txt(slide, 1.5, 3.3, 4, 0.5, "POST /memories", size=18, color=YELLOW, bold=True, font="Consolas")
txt(slide, 1.5, 3.9, 4, 2.8, """Store with auto-classification
→ Extract entities & links
→ Classify: fact, preference...
→ Encrypt at rest (AES-256)
→ Vectorize for semantic search""", size=13, color=WHITE)

add_card(slide, 6, 3, 4.5, 4.5, DARK_GRAY, YELLOW)
txt(slide, 6.3, 3.3, 4, 0.5, "POST /search", size=18, color=YELLOW, bold=True, font="Consolas")
txt(slide, 6.3, 3.9, 4, 2.8, """Retrieve relevant memories
→ Semantic + keyword hybrid
→ Importance-weighted ranking
→ Time-aware filtering
→ Namespace isolation""", size=13, color=WHITE)

add_card(slide, 10.8, 3, 4.5, 4.5, DARK_GRAY, YELLOW)
txt(slide, 11.1, 3.3, 4, 0.5, "POST /context", size=18, color=YELLOW, bold=True, font="Consolas")
txt(slide, 11.1, 3.9, 4, 2.8, """Session context injection
→ Critical + high tier memories
→ Token-optimized format
→ Mid-conversation refresh
→ Auto-promote on access""", size=13, color=WHITE)

mono(slide, 1.2, 7.8, 14, 0.6, """# Add memory to any agent in 2 lines
context = engrm.context()  # Get relevant memories
engrm.store("User prefers dark mode", type="preference")""", size=12, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "04", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 5: Intelligence Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Memory that thinks", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 10, 0.6, "Built-in intelligence. Not just storage.", size=20, color=DARK_GRAY)

features = [
    ("Tiered Importance", "Critical → Working → High → Normal\nAuto-promotion: 5+ accesses → high, 15+ → critical"),
    ("Reflect & Synthesize", "LLM consolidates recent memories into insights.\nRemove redundancy, strengthen patterns."),
    ("Smart Decay", "Ebbinghaus curves (×0.95/week default).\nImportant = slow decay. Reinforced = stronger."),
    ("Entity Linking", "Auto-extract people, places, concepts.\nsame_entity edges connect related memories."),
]
for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = 1.2 + col * 7.3
    y = 3.3 + row * 2.5
    add_card(slide, x, y, 6.8, 2.2, BLACK)
    txt(slide, x + 0.3, y + 0.25, 6, 0.5, title, size=17, color=YELLOW, bold=True)
    txt(slide, x + 0.3, y + 0.85, 6.2, 1.3, desc, size=12, color=WHITE)

mono(slide, 14.5, 8.2, 1.2, 0.4, "05", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 6: Security - User & Agent Only
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "User & Agent only", size=48, color=WHITE, bold=True)
txt(slide, 1.2, 2.6, 10, 0.8, "Core principle: Nobody else reads your data. Not Engrm. Not the database.", size=20, color=GRAY)

add_card(slide, 1.2, 4, 4.3, 2.6, DARK_GRAY, YELLOW)
txt(slide, 1.5, 4.3, 4, 0.5, "Agent sends", size=16, color=YELLOW, bold=True)
txt(slide, 1.5, 4.9, 3.8, 1.6, "Plaintext via API.\nSimple integration.\nNo crypto on client.", size=14, color=WHITE)

txt(slide, 5.8, 5.1, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 6.5, 4, 4.3, 2.6, DARK_GRAY)
txt(slide, 6.8, 4.3, 4, 0.5, "Server encrypts", size=16, color=WHITE, bold=True)
txt(slide, 6.8, 4.9, 3.8, 1.6, "AES-256-GCM.\nPer-user derived keys.\nEncrypted before storage.", size=14, color=GRAY)

txt(slide, 11.1, 5.1, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 11.5, 4, 4.3, 2.6, DARK_GRAY, YELLOW)
txt(slide, 11.8, 4.3, 4, 0.5, "DB sees only", size=16, color=YELLOW, bold=True)
txt(slide, 11.8, 4.9, 3.8, 1.6, "Ciphertext.\nCan't decrypt.\nBreach = useless bytes.", size=14, color=WHITE)

txt(slide, 1.2, 7.2, 14, 0.5, "Encryption is non-negotiable. Removed for performance = immediately reverted.", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "06", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 7: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Why we win", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.5, 12, 0.6, "Cognee raised $7.5M validating the category. We're simpler, faster, privacy-first.", size=20, color=DARK_GRAY)

add_card(slide, 1.2, 3.5, 13.5, 4.2, BLACK)

# Headers
txt(slide, 1.5, 3.7, 3.5, 0.4, "", size=14, color=GRAY)
txt(slide, 5.2, 3.7, 3.5, 0.4, "Cognee", size=14, color=GRAY, bold=True)
txt(slide, 9.5, 3.7, 3.5, 0.4, "Mem0", size=14, color=GRAY, bold=True)
txt(slide, 12.5, 3.7, 2, 0.4, "Engrm", size=14, color=YELLOW, bold=True)

rows = [
    ("Setup", "Days", "Minutes", "Minutes"),
    ("Integration", "ECL pipeline", "Simple API", "2 endpoints"),
    ("Infrastructure", "Graph + Vector + SQL", "Hosted", "Managed"),
    ("Encryption", "Enterprise tier", "Optional", "Default ✓"),
    ("Self-improving", "Complex tuning", "Limited", "Built-in ✓"),
    ("Privacy model", "Standard", "Standard", "User-only ✓"),
]
for i, (label, cognee, mem0, engrm) in enumerate(rows):
    y = 4.2 + i * 0.55
    txt(slide, 1.5, y, 3.5, 0.5, label, size=12, color=GRAY)
    txt(slide, 5.2, y, 3.5, 0.5, cognee, size=12, color=WHITE)
    txt(slide, 9.5, y, 2.5, 0.5, mem0, size=12, color=WHITE)
    txt(slide, 12.5, y, 2, 0.5, engrm, size=12, color=YELLOW)

mono(slide, 14.5, 8.2, 1.2, 0.4, "07", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 8: DX & SDK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Developer experience", size=48, color=WHITE, bold=True)
txt(slide, 1.2, 2.6, 10, 0.6, "Memory should feel automatic. Integration in minutes.", size=20, color=GRAY)

add_card(slide, 1.2, 3.5, 7, 4, DARK_GRAY, YELLOW)
txt(slide, 1.5, 3.7, 6, 0.5, "Python SDK", size=18, color=YELLOW, bold=True)
mono(slide, 1.5, 4.3, 6.5, 3, """from engrm import Engrm

client = Engrm(api_key="mem_...")

# Store automatically
client.store("User loves hiking")

# Recall with context
memories = client.search("outdoor activities")

# Inject into prompt
context = client.context()""", size=11, color=WHITE)

add_card(slide, 8.7, 3.5, 6.3, 4, DARK_GRAY)
txt(slide, 9, 3.7, 5.5, 0.5, "What you get", size=18, color=WHITE, bold=True)
txt(slide, 9, 4.3, 5.8, 3, """✓ Auto memory classification
✓ Entity extraction included
✓ Semantic + keyword search
✓ Token-optimized context
✓ Memory Browser dashboard
✓ Usage analytics
✓ MCP server for Claude
✓ REST API + TypeScript SDK""", size=13, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "08", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 9: Status & Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Status & roadmap", size=48, color=BLACK, bold=True)

add_card(slide, 1.2, 3, 7, 4.5, BLACK)
txt(slide, 1.5, 3.2, 6, 0.5, "Live in production", size=18, color=YELLOW, bold=True)
txt(slide, 1.5, 3.8, 6.5, 3.5, """✓ Core API (store, search, context)
✓ Server-side encryption (AES-256-GCM)
✓ Memory types & classification
✓ Tiered importance system
✓ Reflect, consolidate, decay
✓ Entity extraction & linking
✓ Memory Browser dashboard
✓ Python SDK + MCP server
✓ LongMemEval benchmark integration""", size=13, color=WHITE)

add_card(slide, 8.7, 3, 6, 4.5, DARK_GRAY)
txt(slide, 9, 3.2, 5.5, 0.5, "Next 90 days", size=18, color=WHITE, bold=True)
txt(slide, 9, 3.8, 5.5, 3.5, """→ LangChain integration
→ OpenAI Agents SDK adapter
→ Proactive recall (auto-surface)
→ Contradiction detection
→ Multi-agent memory sharing
→ Enterprise SSO + audit logs
→ On-prem deployment option""", size=13, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "09", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 10: Pricing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Simple pricing", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.5, "Start free. Scale with usage.", size=20, color=DARK_GRAY)

tiers = [
    ("FREE", "$0", "5K memories/mo\n100MB storage\nCore features\nCommunity support"),
    ("PRO", "$29", "50K memories/mo\n1GB storage\nReflect & analytics\nPriority support"),
    ("ENTERPRISE", "Custom", "Unlimited scale\nOn-prem option\nSSO + SLA\nDedicated support"),
]
for i, (tier, price, features) in enumerate(tiers):
    x = 1.2 + i * 5
    is_pro = tier == "PRO"
    add_card(slide, x, 3.5, 4.5, 4.2, BLACK if is_pro else DARK_GRAY)
    txt(slide, x + 0.3, 3.8, 4, 0.4, tier, size=12, color=YELLOW if is_pro else GRAY, bold=True, font="Consolas")
    txt(slide, x + 0.3, 4.3, 4, 0.7, price, size=44, color=YELLOW if is_pro else WHITE, bold=True)
    txt(slide, x + 0.3, 5.2, 4, 2.2, features, size=13, color=WHITE if is_pro else GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "10", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 11: CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)
add_shape(slide, 9, 2, 7, 7, DARK_GRAY)

txt(slide, 1.2, 2.8, 8, 1.6, "Give your agents\na memory", size=64, color=WHITE, bold=True)
txt(slide, 1.2, 5.3, 7, 1.2, "The cognitive layer for AI agents.\nSelf-improving. Encrypted. Fast to integrate.", size=22, color=GRAY)

add_card(slide, 1.2, 7, 3.5, 0.7, YELLOW)
txt(slide, 1.4, 7.1, 3.2, 0.5, "Start building →", size=16, color=BLACK, bold=True)

txt(slide, 1.2, 8, 4, 0.4, "engrm.xyz", size=18, color=YELLOW, bold=True, font="Consolas")

mono(slide, 14.5, 8.2, 1.2, 0.4, "11", color=YELLOW, align=PP_ALIGN.RIGHT)

# Save
output_path = '/Users/clawdaddy/clawd/projects/engrm/Engrm-Pitch-Deck-v3.pptx'
prs.save(output_path)
print(f"Created: {output_path}")
