#!/usr/bin/env python3
"""
Engrm Pitch Deck v2 - a16z speedrun
Updated with latest features: encryption-first, self-improving memory, competitive validation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Yellow/Black palette (a16z style)
YELLOW = RGBColor(255, 219, 21)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(30, 30, 30)
GRAY = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)
CYAN = RGBColor(34, 211, 238)

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

def add_yellow_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = YELLOW
    bg.line.fill.background()

def add_black_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK
    bg.line.fill.background()

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, fill_color=DARK_GRAY, border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(2)
    else:
        card.line.fill.background()
    return card

def txt(slide, left, top, width, height, text, size=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align

def mono(slide, left, top, width, height, text, size=11, color=GRAY, align=PP_ALIGN.LEFT):
    txt(slide, left, top, width, height, text, size, color, False, align, "Consolas")

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)
add_shape(slide, 10, 1.5, 6, 6, DARK_GRAY)

txt(slide, 1.2, 3, 9, 1.6, "Memory for\nAI agents", size=72, color=BLACK, bold=True)
txt(slide, 1.2, 5.5, 8, 1, "Self-improving memory that compounds over time.\nEncrypted by default. Works with any LLM.", size=22, color=DARK_GRAY)
txt(slide, 1.2, 7.2, 6, 0.5, "engrm.xyz", size=18, color=BLACK, bold=True, font="Consolas")
mono(slide, 14.5, 8.2, 1.2, 0.4, "01", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 2: Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1.2, "Every AI agent\nwakes up with amnesia", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 4, 8, 1.5, "ChatGPT forgets you between sessions.\nCustom agents lose context on restart.\nRAG retrieves documents, not relationships.", size=22, color=DARK_GRAY)

add_card(slide, 10, 2, 4.5, 1.8, BLACK)
txt(slide, 10.3, 2.2, 4, 0.9, "0", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 3, 4, 0.5, "MEMORIES RETAINED\nBETWEEN SESSIONS", size=10, color=GRAY)

add_card(slide, 10, 4, 4.5, 1.8, BLACK)
txt(slide, 10.3, 4.2, 4, 0.9, "$7.5M", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 5, 4, 0.5, "COGNEE SEED ROUND\nVALIDATES CATEGORY", size=10, color=GRAY)

add_card(slide, 10, 6, 4.5, 1.8, BLACK)
txt(slide, 10.3, 6.2, 4, 0.9, "$47B", size=56, color=YELLOW, bold=True)
mono(slide, 10.3, 7, 4, 0.5, "AI AGENT MARKET\nBY 2030", size=10, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "02", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 3: Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.8, 10, 1, "A brain, not a database", size=52, color=BLACK, bold=True)
txt(slide, 1.2, 3, 8, 0.8, "Memory that works like human cognition — and improves itself.", size=22, color=DARK_GRAY)

pillars = [
    ("🧠", "Biological decay", "Important memories strengthen.\nStale ones fade naturally."),
    ("🔗", "Self-improving", "Reflect, consolidate, compress.\nGets smarter over time."),
    ("🔒", "Encrypted by default", "Client-side AES-256.\nWe can't read your data."),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 1.2 + i * 5
    add_card(slide, x, 4, 4.5, 3.2, BLACK)
    txt(slide, x + 0.2, 4.3, 4, 0.6, icon, size=36, color=WHITE)
    txt(slide, x + 0.2, 5, 4, 0.5, title, size=20, color=YELLOW, bold=True)
    txt(slide, x + 0.2, 5.6, 4, 1.4, desc, size=14, color=WHITE)

mono(slide, 14.5, 8.2, 1.2, 0.4, "03", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 4: How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 12, 1, "Two endpoints. That's it.", size=48, color=WHITE, bold=True)

add_card(slide, 1.2, 3, 6.5, 4.5, DARK_GRAY, YELLOW)
txt(slide, 1.5, 3.3, 6, 0.5, "POST /store", size=20, color=YELLOW, bold=True, font="Consolas")
txt(slide, 1.5, 4, 6, 2.5, """Agent saves memory (encrypted)
→ Auto-extract entities
→ Vectorize for semantic search
→ Link to related memories
→ Apply decay curves""", size=14, color=WHITE)

add_card(slide, 8.3, 3, 6.5, 4.5, DARK_GRAY, YELLOW)
txt(slide, 8.6, 3.3, 6, 0.5, "POST /recall", size=20, color=YELLOW, bold=True, font="Consolas")
txt(slide, 8.6, 4, 6, 2.5, """Query relevant memories
→ Semantic + keyword search
→ Boost by feedback score
→ Filter by time/namespace
→ Inject into LLM context""", size=14, color=WHITE)

mono(slide, 1.2, 7.8, 14, 0.8, """# That's it. Two lines to add memory to any agent.
memories = engrm.recall("user preferences")
engrm.store("User prefers dark mode")""", size=13, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "04", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 5: Self-Improving Memory (NEW)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Memory that improves itself", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.6, "Built-in LLM-powered maintenance. Not just storage.", size=20, color=DARK_GRAY)

features = [
    ("Reflect", "Consolidate recent memories into insights.\nRemove redundant facts, strengthen patterns."),
    ("Feedback", "Rate recall results → boost relevance.\nMemories that help get reinforced."),
    ("Compact", "Auto-merge similar memories.\nKeep signal, reduce noise."),
    ("Decay", "Ebbinghaus curves. Important = slow decay.\nAccessed often = stronger retention."),
]
for i, (title, desc) in enumerate(features):
    row = i // 2
    col = i % 2
    x = 1.2 + col * 7.3
    y = 3.5 + row * 2.4
    add_card(slide, x, y, 6.8, 2.1, BLACK)
    txt(slide, x + 0.3, y + 0.2, 6, 0.5, title, size=18, color=YELLOW, bold=True)
    txt(slide, x + 0.3, y + 0.8, 6, 1.2, desc, size=13, color=WHITE)

mono(slide, 14.5, 8.2, 1.2, 0.4, "05", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 6: Security (Encryption-First)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Encrypted by default", size=48, color=WHITE, bold=True)
txt(slide, 1.2, 2.6, 8, 0.8, "Not optional. Not an add-on. The foundation.", size=20, color=GRAY)

add_card(slide, 1.2, 4, 4, 2.5, DARK_GRAY, YELLOW)
txt(slide, 1.5, 4.3, 3.5, 0.5, "Client encrypts", size=16, color=YELLOW, bold=True)
txt(slide, 1.5, 4.9, 3.5, 1.5, "AES-256-GCM\nbefore upload.\nWe never see plaintext.", size=14, color=WHITE)

txt(slide, 5.5, 5, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 6.5, 4, 4, 2.5, DARK_GRAY)
txt(slide, 6.8, 4.3, 3.5, 0.5, "We store", size=16, color=GRAY, bold=True)
txt(slide, 6.8, 4.9, 3.5, 1.5, "Ciphertext only.\nCan't decrypt.\nCan't read.", size=14, color=GRAY)

txt(slide, 10.8, 5, 1, 0.5, "→", size=28, color=YELLOW)

add_card(slide, 11.5, 4, 4, 2.5, DARK_GRAY, YELLOW)
txt(slide, 11.8, 4.3, 3.5, 0.5, "Client decrypts", size=16, color=YELLOW, bold=True)
txt(slide, 11.8, 4.9, 3.5, 1.5, "On retrieval.\nYour keys.\nYour data.", size=14, color=WHITE)

txt(slide, 1.2, 7.2, 12, 0.5, "API rejects plaintext. Vault setup required before creating keys.", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "06", color=YELLOW, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 7: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Why we win", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.5, 12, 0.6, "Cognee raised $7.5M validating the category. We're the simpler, faster alternative.", size=20, color=DARK_GRAY)

# Comparison table
add_card(slide, 1.2, 3.5, 13.5, 4, BLACK)

# Headers
txt(slide, 1.5, 3.7, 4, 0.4, "", size=14, color=GRAY)
txt(slide, 5.5, 3.7, 4, 0.4, "Cognee", size=14, color=GRAY, bold=True)
txt(slide, 10, 3.7, 4, 0.4, "Engrm", size=14, color=YELLOW, bold=True)

rows = [
    ("Setup time", "Days/weeks", "5 minutes"),
    ("Integration", "ECL pipeline, 38 connectors", "2 endpoints"),
    ("Infrastructure", "Relational + Vector + Graph", "We manage it"),
    ("Encryption", "Enterprise add-on", "Default, enforced"),
    ("Self-improving", "Knowledge graph tuning", "Built-in reflect/compact"),
]
for i, (label, cognee, engrm) in enumerate(rows):
    y = 4.2 + i * 0.6
    txt(slide, 1.5, y, 4, 0.5, label, size=13, color=GRAY)
    txt(slide, 5.5, y, 4, 0.5, cognee, size=13, color=WHITE)
    txt(slide, 10, y, 4, 0.5, engrm, size=13, color=YELLOW)

mono(slide, 14.5, 8.2, 1.2, 0.4, "07", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 8: Traction / Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Status & roadmap", size=48, color=BLACK, bold=True)

add_card(slide, 1.2, 3, 7, 4.5, BLACK)
txt(slide, 1.5, 3.2, 6, 0.5, "Live today", size=18, color=YELLOW, bold=True)
txt(slide, 1.5, 3.8, 6.5, 3.5, """✓ Core API (store, recall, search)
✓ Encryption enforced by default
✓ Feedback loops & decay curves
✓ Self-improvement (reflect, compact)
✓ Entity extraction
✓ PDF ingestion
✓ MCP server for Claude
✓ REST API + Python SDK""", size=14, color=WHITE)

add_card(slide, 8.7, 3, 6, 4.5, DARK_GRAY)
txt(slide, 9, 3.2, 5.5, 0.5, "Next 90 days", size=18, color=WHITE, bold=True)
txt(slide, 9, 3.8, 5.5, 3.5, """→ LangChain integration
→ OpenAI Agents SDK
→ Proactive recall (auto-surface)
→ Contradiction detection
→ Multi-agent memory sharing
→ Enterprise SSO
→ On-prem option""", size=14, color=GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "08", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 9: Pricing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_yellow_bg(slide)

txt(slide, 1.2, 1.5, 10, 1, "Simple pricing", size=48, color=BLACK, bold=True)
txt(slide, 1.2, 2.6, 8, 0.5, "Start free. Scale with usage.", size=20, color=DARK_GRAY)

tiers = [
    ("FREE", "$0", "5K memories/mo\n100MB storage\nCore features"),
    ("PRO", "$29", "50K memories/mo\n1GB storage\nReflect & compact"),
    ("ENTERPRISE", "Custom", "Unlimited\nOn-prem option\nSSO + SLA"),
]
for i, (tier, price, features) in enumerate(tiers):
    x = 1.2 + i * 5
    is_pro = tier == "PRO"
    add_card(slide, x, 3.5, 4.5, 4, BLACK if is_pro else DARK_GRAY)
    txt(slide, x + 0.3, 3.8, 4, 0.4, tier, size=12, color=YELLOW if is_pro else GRAY, bold=True, font="Consolas")
    txt(slide, x + 0.3, 4.3, 4, 0.7, price, size=44, color=YELLOW if is_pro else WHITE, bold=True)
    txt(slide, x + 0.3, 5.2, 4, 2, features, size=14, color=WHITE if is_pro else GRAY)

mono(slide, 14.5, 8.2, 1.2, 0.4, "09", color=GRAY, align=PP_ALIGN.RIGHT)

# ============================================================
# SLIDE 10: CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_black_bg(slide)
add_shape(slide, 9, 2, 7, 7, DARK_GRAY)

txt(slide, 1.2, 3, 8, 1.6, "Give your agents\na memory", size=64, color=WHITE, bold=True)
txt(slide, 1.2, 5.5, 7, 1, "Self-improving. Encrypted. Fast to integrate.\nThe memory layer AI agents deserve.", size=22, color=GRAY)

add_card(slide, 1.2, 7, 3.5, 0.7, YELLOW)
txt(slide, 1.4, 7.1, 3.2, 0.5, "Start building →", size=16, color=BLACK, bold=True)

txt(slide, 1.2, 8, 4, 0.4, "engrm.xyz", size=18, color=YELLOW, bold=True, font="Consolas")

mono(slide, 14.5, 8.2, 1.2, 0.4, "10", color=YELLOW, align=PP_ALIGN.RIGHT)

# Save
prs.save('/Users/clawdaddy/clawd/projects/engrm/Engrm-Pitch-Deck-v2.pptx')
print("Created: Engrm-Pitch-Deck-v2.pptx")
